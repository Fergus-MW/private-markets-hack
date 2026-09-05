"""Optional local conversion checks for legacy Office, ODT, RTF, and EPUB."""
import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory
import httpx
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

with TemporaryDirectory() as directory, httpx.Client(timeout=240) as client:
    root = Path(directory)
    doc = Document(); doc.add_paragraph("Investment review: revenue was GBP 100 million.")
    doc.save(root / "report.docx")
    workbook = Workbook(); workbook.active.append(["Revenue", 100]); workbook.save(root / "accounts.xlsx")
    deck = Presentation(); slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Investment review"; slide.placeholders[1].text = "Revenue was GBP 100 million."
    deck.save(root / "review.pptx")
    for source, extension in [("report.docx", "doc"), ("report.docx", "rtf"),
                              ("report.docx", "odt"), ("accounts.xlsx", "xls"), ("review.pptx", "ppt")]:
        subprocess.run(["libreoffice", "--headless", "--convert-to", extension, "--outdir", directory,
                        str(root / source)], check=True, capture_output=True, timeout=60)
        path = (root / source).with_suffix("." + extension)
        assert path.exists(), path
        response = client.post("http://127.0.0.1:8080/documents", files={"file": (path.name, path.read_bytes())})
        assert response.status_code == 200, (extension, response.text)
        assert response.json()["chunks"] > 0
        print("PASS:", extension, flush=True)
    (root / "report.md").write_text("# Investment review\n\nRevenue was GBP 100 million.")
    subprocess.run(["pandoc", str(root / "report.md"), "-o", str(root / "report.epub")], check=True, capture_output=True)
    for filename, data in [("report.epub", (root / "report.epub").read_bytes()),
                           ("report.xml", b"<document><paragraph>Revenue was GBP 100 million.</paragraph></document>"),
                           ("report.rst", b"Investment review\n=================\n\nRevenue was GBP 100 million.")]:
        response = client.post("http://127.0.0.1:8080/documents", files={"file": (filename, data)})
        assert response.status_code == 200, (filename, response.text)
        print("PASS:", filename, flush=True)
