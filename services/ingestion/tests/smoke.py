"""Local container integration test with real Unstructured parsers and SurrealDB."""
import io
import os
import time
import httpx
from app.store import Store

base = "http://127.0.0.1:8080"
text = b"Quarterly investment review. Revenue increased to GBP 100 million. The proposed acquisition remains subject to diligence."

def upload(client, filename, content):
    response = client.post(base + "/documents", files={"file": (filename, content)})
    assert response.status_code == 200, (filename, response.text)
    result = response.json()
    assert result["chunks"] > 0 and "people" not in result and "companies" not in result, result
    context = client.get(base + result["context_url"]).json()
    assert context["chunks"] and context["chunks"][0]["sources"], context
    assert all(len(chunk["text"]) <= 4000 for chunk in context["chunks"])
    assert context["chunks"][0]["citation"]["filename"] == filename
    elements = client.get(base + "/documents/" + result["document_id"] + "/elements").json()
    assert elements["elements"]
    return result, context, elements

with httpx.Client(timeout=240) as client:
    for attempt in range(30):
        try:
            if client.get(base + "/readyz").status_code == 200:
                break
        except httpx.ConnectError:
            pass
        time.sleep(1)
    else:
        raise AssertionError("Service did not become ready")
    result, context, _ = upload(client, "sample.txt", text)
    assert upload(client, "sample.txt", text)[0] == result
    assert client.get(base + "/documents/" + result["document_id"]).json()["chunk_count"] == result["chunks"]
    assert client.get(base + result["context_url"] + "?max_characters=1").status_code == 422
    assert client.get(base + "/documents/" + "0" * 64 + "/context").status_code == 404
    assert client.post(base + "/documents", files={"file": ("bad.exe", text)}).status_code == 415
    assert client.post(base + "/documents", files={"file": ("empty.txt", b"")}).status_code == 422
    assert client.post(base + "/documents", files={"file": ("broken.pdf", b"not a PDF")}).status_code == 422
    assert client.post(base + "/documents", files={"file": ("huge.txt", b"x" * (20 * 1024 * 1024 + 1))}).status_code == 413

    from docx import Document
    doc = Document()
    doc.add_heading("Investment review", 1)
    doc.add_paragraph(text.decode())
    buffer = io.BytesIO(); doc.save(buffer)
    upload(client, "sample.docx", buffer.getvalue())

    from pptx import Presentation
    deck = Presentation(); slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Investment review"
    slide.placeholders[1].text = text.decode()
    buffer = io.BytesIO(); deck.save(buffer)
    _, slides, _ = upload(client, "sample.pptx", buffer.getvalue())
    assert any(s.get("page_number") == 1 for c in slides["chunks"] for s in c["sources"])

    from openpyxl import Workbook
    book = Workbook(); sheet = book.active; sheet.title = "Financials"
    sheet.append(["Metric", "Value"]); sheet.append(["Revenue", 100]); sheet.append(["EBITDA", 20])
    buffer = io.BytesIO(); book.save(buffer)
    _, sheets, tables = upload(client, "sample.xlsx", buffer.getvalue())
    assert any(e["metadata"].get("text_as_html") for e in tables["elements"]), tables
    assert any(s.get("page_name") == "Financials" for c in sheets["chunks"] for s in c["sources"]), sheets
    upload(client, "sample.csv", b"Metric,Value\nRevenue,100\nEBITDA,20\n")
    upload(client, "sample.tsv", b"Metric\tValue\nRevenue\t100\nEBITDA\t20\n")
    upload(client, "sample.md", b"# Investment review\n\n" + text)
    upload(client, "sample.html", b"<h1>Investment review</h1><p>" + text + b"</p><table><tr><td>Revenue</td><td>100</td></tr></table>")
    upload(client, "sample.eml", b"From: analyst@example.com\nTo: team@example.com\nSubject: Investment review\nMIME-Version: 1.0\nContent-Type: text/plain; charset=utf-8\n\n" + text)

    from PIL import Image, ImageDraw, ImageFont
    image = Image.new("RGB", (2200, 500), "white")
    ImageDraw.Draw(image).text((50, 100), text.decode(), fill="black", font=ImageFont.load_default(size=28))
    buffer = io.BytesIO(); image.save(buffer, format="PDF", resolution=150)
    upload(client, "scanned.pdf", buffer.getvalue())
    buffer = io.BytesIO(); image.save(buffer, format="PNG")
    upload(client, "scanned.png", buffer.getvalue())

    # Oversized source element must split without losing citation or exceeding budget.
    long_text = b"Investment diligence identifies operational risks and growth opportunities. " * 160
    long_doc, _, _ = upload(client, "long.txt", long_text)
    offset, chunk_ids = 0, []
    while offset is not None:
        page = client.get(base + long_doc["context_url"], params={"offset": offset, "max_characters": 4000}).json()
        assert page["characters"] <= 4000
        assert page["chunks"], page
        chunk_ids.extend(c["chunk_id"] for c in page["chunks"])
        offset = page["next_offset"]
    assert len(chunk_ids) == long_doc["chunks"] and len(set(chunk_ids)) == len(chunk_ids)

Store().query("DEFINE USER OVERWRITE ingestion ON DATABASE PASSWORD 'localtestonly' ROLES EDITOR;")
os.environ.update(SURREAL_AUTH_LEVEL="database", SURREAL_USER="ingestion", SURREAL_PASSWORD="localtestonly")
doc = Store().get(result["document_id"])
Store().save(doc)
assert Store().get(result["document_id"])["chunks"] == doc["chunks"]
print("PASS: document formats, OCR, table/slide/sheet provenance, context pagination, retries, validation, and database-scoped writes")
